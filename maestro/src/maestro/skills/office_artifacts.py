"""Safe, structured creation and delivery of Word/PowerPoint artifacts."""

from __future__ import annotations

import json
import mimetypes
import re
import uuid
from pathlib import Path
from urllib.parse import quote

from maestro.skills.schemas import SkillValidationError


_SAFE_FILENAME = re.compile(r"[^\w\-.\u4e00-\u9fff]+", re.UNICODE)


def _filename(value: str, suffix: str) -> str:
    name = Path(value or f"artifact{suffix}").name
    if not name.lower().endswith(suffix):
        name += suffix
    name = _SAFE_FILENAME.sub("_", name).strip("._")
    if not name:
        name = f"artifact{suffix}"
    return name[:180]


def artifact_metadata(path: Path, root: Path) -> dict:
    relative = path.resolve().relative_to(root.resolve())
    parts = [quote(part, safe="") for part in relative.parts]
    return {
        "name": path.name,
        "size": path.stat().st_size,
        "content_type": mimetypes.guess_type(path.name)[0] or "application/octet-stream",
        "download_url": "/artifacts/" + "/".join(parts),
    }


class OfficeArtifactService:
    """Generate bounded Office files from content specs; never execute model code."""

    def __init__(self, output_dir: Path):
        self.output_dir = Path(output_dir) / "artifacts"

    def create(self, params: dict) -> dict:
        kind = str(params.get("kind", "")).lower()
        if kind not in {"docx", "pptx"}:
            raise SkillValidationError("kind 必须是 docx 或 pptx")
        title = str(params.get("title", "")).strip()
        if not title or len(title) > 200:
            raise SkillValidationError("标题不能为空且不能超过 200 字")
        subtitle = str(params.get("subtitle", "")).strip()[:300]
        sections = params.get("sections", [])
        if not isinstance(sections, list) or not 1 <= len(sections) <= 30:
            raise SkillValidationError("sections 必须包含 1~30 个章节/页面")
        normalized = [self._section(item) for item in sections]

        run_id = f"office-{uuid.uuid4().hex[:12]}"
        run_dir = self.output_dir / run_id
        run_dir.mkdir(parents=True, exist_ok=False)
        suffix = f".{kind}"
        path = run_dir / _filename(str(params.get("filename", title)), suffix)
        if kind == "docx":
            self._create_docx(path, title, subtitle, normalized)
        else:
            self._create_pptx(path, title, subtitle, normalized)
        return {
            "status": "completed",
            "kind": kind,
            "artifact": artifact_metadata(path, self.output_dir),
        }

    @staticmethod
    def _section(item: object) -> dict:
        if not isinstance(item, dict):
            raise SkillValidationError("每个 section 必须是对象")
        title = str(item.get("title", "")).strip()[:160]
        paragraphs = item.get("paragraphs", [])
        bullets = item.get("bullets", [])
        if not isinstance(paragraphs, list) or not isinstance(bullets, list):
            raise SkillValidationError("paragraphs 与 bullets 必须是字符串数组")
        clean_paragraphs = [str(value).strip()[:3000] for value in paragraphs if str(value).strip()]
        clean_bullets = [str(value).strip()[:800] for value in bullets if str(value).strip()]
        if len(clean_paragraphs) > 20 or len(clean_bullets) > 20:
            raise SkillValidationError("单个 section 的段落或要点不能超过 20 条")
        if not title and not clean_paragraphs and not clean_bullets:
            raise SkillValidationError("section 不能全空")
        return {"title": title, "paragraphs": clean_paragraphs, "bullets": clean_bullets}

    @staticmethod
    def _create_docx(path: Path, title: str, subtitle: str, sections: list[dict]) -> None:
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            from docx.shared import Inches, Pt, RGBColor
        except ImportError as exc:  # pragma: no cover - packaging/configuration failure
            raise SkillValidationError("后端缺少 python-docx，无法生成 Word 文件") from exc

        doc = Document()
        sec = doc.sections[0]
        sec.top_margin = Inches(0.8)
        sec.bottom_margin = Inches(0.75)
        sec.left_margin = Inches(0.9)
        sec.right_margin = Inches(0.9)
        styles = doc.styles
        styles["Normal"].font.name = "Arial"
        styles["Normal"].font.size = Pt(11)
        styles["Heading 1"].font.name = "Arial"
        styles["Heading 1"].font.size = Pt(18)
        styles["Heading 1"].font.color.rgb = RGBColor(31, 78, 121)

        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(6)
        run = p.add_run(title)
        run.bold = True
        run.font.name = "Arial"
        run.font.size = Pt(26)
        run.font.color.rgb = RGBColor(24, 55, 90)
        if subtitle:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_after = Pt(24)
            run = p.add_run(subtitle)
            run.font.name = "Arial"
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(90, 101, 115)

        for index, item in enumerate(sections):
            if item["title"]:
                heading = doc.add_heading(item["title"], level=1)
                heading.paragraph_format.keep_with_next = True
            for text in item["paragraphs"]:
                p = doc.add_paragraph(text)
                p.paragraph_format.space_after = Pt(7)
                p.paragraph_format.line_spacing = 1.15
            for text in item["bullets"]:
                p = doc.add_paragraph(text, style="List Bullet")
                p.paragraph_format.space_after = Pt(4)
            if index < len(sections) - 1:
                doc.add_paragraph().paragraph_format.space_after = Pt(2)

        footer = sec.footer.paragraphs[0]
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        field = OxmlElement("w:fldSimple")
        field.set(qn("w:instr"), "PAGE")
        footer._p.append(field)
        doc.save(path)

    @staticmethod
    def _create_pptx(path: Path, title: str, subtitle: str, sections: list[dict]) -> None:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as exc:  # pragma: no cover - packaging/configuration failure
            raise SkillValidationError("后端缺少 python-pptx，无法生成 PowerPoint 文件") from exc

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        navy, blue, pale, ink = RGBColor(20, 42, 72), RGBColor(39, 110, 184), RGBColor(232, 241, 250), RGBColor(32, 42, 55)

        cover = prs.slides.add_slide(prs.slide_layouts[6])
        cover.background.fill.solid()
        cover.background.fill.fore_color.rgb = navy
        box = cover.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.5))
        para = box.text_frame.paragraphs[0]
        para.text = title
        para.font.name = "Arial"
        para.font.size = Pt(34)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
        if subtitle:
            sub = cover.shapes.add_textbox(Inches(1.4), Inches(3.65), Inches(10.5), Inches(0.7))
            para = sub.text_frame.paragraphs[0]
            para.text = subtitle
            para.font.name = "Arial"
            para.font.size = Pt(18)
            para.font.color.rgb = pale
            para.alignment = PP_ALIGN.CENTER

        for page_number, item in enumerate(sections, start=2):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
            accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
            accent.fill.solid(); accent.fill.fore_color.rgb = blue; accent.line.fill.background()
            head = slide.shapes.add_textbox(Inches(0.85), Inches(0.55), Inches(11.6), Inches(0.7))
            p = head.text_frame.paragraphs[0]
            p.text = item["title"] or f"第 {page_number - 1} 部分"
            p.font.name = "Arial"; p.font.size = Pt(27); p.font.bold = True; p.font.color.rgb = navy
            body = slide.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.2), Inches(4.95))
            tf = body.text_frame; tf.word_wrap = True; tf.clear()
            lines = [(value, False) for value in item["paragraphs"]] + [(value, True) for value in item["bullets"]]
            if not lines:
                lines = [("内容待补充", False)]
            for idx, (value, bullet) in enumerate(lines):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = value; p.font.name = "Arial"; p.font.size = Pt(18); p.font.color.rgb = ink
                p.space_after = Pt(12); p.level = 0
                if bullet:
                    p.text = f"• {value}"
            marker = slide.shapes.add_textbox(Inches(11.8), Inches(6.85), Inches(0.7), Inches(0.3))
            p = marker.text_frame.paragraphs[0]; p.text = str(page_number); p.font.size = Pt(10); p.font.color.rgb = blue; p.alignment = PP_ALIGN.RIGHT
        prs.save(path)


def format_office_result_markdown(result: dict) -> str:
    artifact = result.get("artifact") or {}
    if result.get("status") != "completed" or not artifact.get("download_url"):
        return f"⚠️ Office 文件生成失败\n\n```json\n{json.dumps(result, ensure_ascii=False)}\n```"
    name = artifact.get("name", "下载文件")
    size = int(artifact.get("size", 0))
    return f"✅ 文件已生成（{size / 1024:.1f} KB）\n\n[下载 {name}]({artifact['download_url']})"
