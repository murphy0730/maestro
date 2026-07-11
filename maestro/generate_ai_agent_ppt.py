#!/usr/bin/env python3
"""生成「AI Agent 深度介绍」PPT — 16:9 宽屏，深色科技风"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色 ──────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)   # 深空背景
CARD_BG      = RGBColor(0x16, 0x1B, 0x22)   # 卡片底色
ACCENT_BLUE  = RGBColor(0x58, 0xA6, 0xFF)   # 科技蓝
ACCENT_GREEN = RGBColor(0x7E, 0xE7, 0x87)   # 翠绿
WHITE        = RGBColor(0xE6, 0xED, 0xF3)   # 主文字
GRAY         = RGBColor(0x8B, 0x94, 0x9E)   # 次要文字
ORANGE       = RGBColor(0xF0, 0x88, 0x3E)   # 警示/高亮
PURPLE       = RGBColor(0xBC, 0x8C, 0xFF)   # 紫色点缀

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── 工具函数 ──────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    """填充纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color=CARD_BG, radius=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        # 小圆角
        shape.adjustments[0] = 0.05
    return shape

def add_circle(slide, left, top, size, color):
    """装饰圆"""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=WHITE, line_spacing=1.3):
    """多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox

def slide_header(slide, number, title, subtitle=""):
    """统一页头: 编号 + 标题"""
    # 顶部装饰线
    add_rect(slide, Inches(0.8), Inches(0.45), Inches(11.7), Pt(2), ACCENT_BLUE)
    # 编号圆圈
    c = add_circle(slide, Inches(0.8), Inches(0.6), Inches(0.45), ACCENT_BLUE)
    c.text_frame.paragraphs[0].text = str(number).zfill(2)
    c.text_frame.paragraphs[0].font.size = Pt(14)
    c.text_frame.paragraphs[0].font.color.rgb = WHITE
    c.text_frame.paragraphs[0].font.bold = True
    c.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 标题
    add_text_box(slide, Inches(1.5), Inches(0.55), Inches(10), Inches(0.5),
                 title, font_size=24, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(0.95), Inches(10), Inches(0.35),
                     subtitle, font_size=13, color=GRAY)

# ═══════════════════════════════════════════════════
#  第1页: 封面
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
add_bg(slide)

# 装饰元素
add_circle(slide, Inches(10.5), Inches(0.8), Inches(2.2),
           RGBColor(0x1A, 0x3A, 0x5C))  # 右上大圆
add_circle(slide, Inches(-0.5), Inches(5.5), Inches(1.5),
           RGBColor(0x0D, 0x2B, 0x1A))  # 左下绿暗圆
add_circle(slide, Inches(7.5), Inches(5.0), Inches(0.6), ACCENT_BLUE)

# 主标题
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "AI Agent", font_size=64, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(2.85), Inches(10), Inches(0.8),
             "智能体深度介绍", font_size=40, bold=True, color=WHITE)

# 副标题
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
             "从原理到应用，全面解读人工智能代理的核心技术与未来趋势",
             font_size=18, color=GRAY)

# 底部装饰线
add_rect(slide, Inches(1.5), Inches(4.8), Inches(3.0), Pt(3), ACCENT_BLUE)

# 日期/作者
add_text_box(slide, Inches(1.5), Inches(6.4), Inches(5), Inches(0.4),
             "2025  ·  技术前沿专题", font_size=14, color=GRAY)

# ═══════════════════════════════════════════════════
#  第2页: 目录
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, 1, "目  录", "CONTENTS")

toc_items = [
    ("01", "什么是 AI Agent", "定义、核心特征与演进历程"),
    ("02", "核心架构", "感知 → 规划 → 执行 → 记忆 四模块"),
    ("03", "关键技术栈", "LLM · Tool Use · RAG · 多Agent协作"),
    ("04", "应用场景矩阵", "六大行业的落地实践"),
    ("05", "主流产品对比", "GPTs · Copilot · Claude · Gemini · 国产方案"),
    ("06", "市场趋势与数据", "规模、增速与五大趋势"),
    ("07", "挑战与伦理", "安全、幻觉、责任归属"),
    ("08", "未来展望", "AGI路线图与终极形态"),
]

for i, (num, title, desc) in enumerate(toc_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.1)
    y = Inches(1.65 + row * 1.3)

    card = add_rect(slide, x, y, Inches(5.6), Inches(1.1), CARD_BG, radius=True)

    # 编号
    add_text_box(slide, Inches(x.inches + 0.3), Inches(y.inches + 0.2),
                 Inches(0.6), Inches(0.5),
                 num, font_size=28, bold=True, color=ACCENT_BLUE)
    # 标题
    add_text_box(slide, Inches(x.inches + 0.95), Inches(y.inches + 0.15),
                 Inches(4.3), Inches(0.4),
                 title, font_size=18, bold=True, color=WHITE)
    # 描述
    add_text_box(slide, Inches(x.inches + 0.95), Inches(y.inches + 0.55),
                 Inches(4.3), Inches(0.35),
                 desc, font_size=12, color=GRAY)

# ═══════════════════════════════════════════════════
#  第3页: 什么是 AI Agent
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, 2, "什么是 AI Agent？", "DEFINITION & EVOLUTION")

# 左侧：定义卡片
card_l = add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2),
                  CARD_BG, radius=True)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(5.0), Inches(0.5),
             "▎核心定义", font_size=20, bold=True, color=ACCENT_BLUE)
add_multiline(slide, Inches(1.2), Inches(2.4), Inches(5.0), Inches(3.8), [
    "AI Agent（人工智能代理）是一种能够",
    "自主感知环境、制定计划、使用工具",
    "并执行复杂任务的智能系统。",
    "",
    "▸ 自主性：无需人类逐步指令",
    "▸ 目标导向：分解目标 → 规划步骤",
    "▸ 工具使用：调用API/数据库/外部程序",
    "▸ 记忆系统：短期上下文 + 长期知识",
    "▸ 反思迭代：根据反馈自我修正",
], font_size=14, color=WHITE, line_spacing=1.4)

# 右侧：时间线
card_r = add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.2),
                  CARD_BG, radius=True)
add_text_box(slide, Inches(7.4), Inches(1.8), Inches(5.0), Inches(0.5),
             "▎演进历程", font_size=20, bold=True, color=ACCENT_GREEN)

timeline = [
    ("2022", "Chain-of-Thought / ReAct 范式提出"),
    ("2023.3", "AutoGPT / BabyAGI 引爆 Agent 概念"),
    ("2023.6", "OpenAI Function Calling 发布"),
    ("2023.11", "GPTs + Assistants API 生态成型"),
    ("2024", "多Agent框架爆发 (CrewAI / AutoGen)"),
    ("2025", "Agent 进入企业级落地深水区"),
]
for i, (year, desc) in enumerate(timeline):
    y = Inches(2.5 + i * 0.65)
    add_circle(slide, Inches(7.5), y, Inches(0.18), ACCENT_GREEN)
    add_text_box(slide, Inches(7.9), Inches(y.inches - 0.05), Inches(1.2), Inches(0.3),
                 year, font_size=12, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(9.2), Inches(y.inches - 0.05), Inches(3.2), Inches(0.3),
                 desc, font_size=12, color=WHITE)

# ═══════════════════════════════════════════════════
#  第4页: 核心架构
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, 3, "核心架构", "ARCHITECTURE: PERCEIVE → PLAN → ACT → MEMORIZE")

modules = [
    ("🔍 感知模块\nPerception", "接收用户输入、读取环境\n上下文、解析多模态数据\n(文本/图像/语音)",
     ACCENT_BLUE),
    ("🧠 规划模块\nPlanning", "任务分解、生成执行计划\n链式推理 (CoT)、\n反思与重规划 (ReAct)",
     ACCENT_GREEN),
    ("🔧 执行模块\nAction", "工具调用 (Function Calling)\n代码执行、API 交互\n浏览器操控 / RPA",
     ORANGE),
    ("💾 记忆模块\nMemory", "短期记忆 (上下文窗口)\n长期记忆 (向量数据库)\n会话持久化与知识积累",
     PURPLE),
]

for i, (title, desc, color) in enumerate(modules):
    x = Inches(0.8 + i * 3.15)
    card = add_rect(slide, x, Inches(1.6), Inches(2.85), Inches(4.8),
                    CARD_BG, radius=True)
    # 顶部色条
    add_rect(slide, Inches(x.inches + 0.3), Inches(1.85), Inches(1.0), Pt(3), color)
    add_multiline(slide, Inches(x.inches + 0.3), Inches(2.1), Inches(2.3), Inches(1.2),
                  title.split("\n"), font_size=16, color=color, line_spacing=1.2)
    add_multiline(slide, Inches(x.inches + 0.3), Inches(3.5), Inches(2.3), Inches(2.5),
                  desc.split("\n"), font_size=12, color=GRAY, line_spacing=1.5)

# 底部流程箭头
for i in range(3):
    ax = Inches(3.65 + i * 3.15)
    add_text_box(slide, ax, Inches(6.6), Inches(0.5), Inches(0.4),
                 "→", font_size=24, bold=True, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
#  第5页: 关键技术栈
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, 4, "关键技术栈", "KEY TECHNOLOGIES")

techs = [
    ("🧠 大语言模型 (LLM)", "GPT-4o / Claude 3.5\nGemini 2.0 / DeepSeek\nQwen / LLaMA 开源", ACCENT_BLUE),
    ("🔧 工具调用", "Function Calling\nMCP 协议 (Anthropic)\nPlugin / Action 生态", ACCENT_GREEN),
    ("📚 RAG 检索增强", "向量数据库 (Pinecone/Milvus)\nEmbedding + 语义检索\nGraph RAG (知识图谱)", ORANGE),
    ("🤝 多Agent协作", "CrewAI / AutoGen\nLangGraph / Swarm\n角色分工 + 消息传递", PURPLE),
    ("🧠 记忆系统", "Mem0 / MemGPT\n短期 + 长期记忆\n用户画像持久化", ACCENT_BLUE),
    ("🛡️ 安全护栏", "Guardrails / NeMo\n输入输出过滤\n权限沙箱机制", RGBColor(0xF7, 0x78, 0x83)),
]

for i, (title, desc, color) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.6 + row * 2.8)

    card = add_rect(slide, x, y, Inches(3.7), Inches(2.5), CARD_BG, radius=True)
    add_rect(slide, Inches(x.inches + 0.3), Inches(y.inches + 0.3), Inches(1.2), Pt(3), color)
    add_text_box(slide, Inches(x.inches + 0.3), Inches(y.inches + 0.5),
                 Inches(3.1), Inches(0.5),
                 title, font_size=15, bold=True, color=color)
    add_multiline(slide, Inches(x.inches + 0.3), Inches(y.inches + 1.1),
                  Inches(3.1), Inches(1.3),
                  desc.split("\n"), font_size=12, color=GRAY, line_spacing=1.5)

# ═══════════════════════════════════════════════════
#  第6页: 应用场景矩阵
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, 5, "应用场景矩阵", "INDUSTRY APPLICATIONS")

scenes = [
    ("🏥 医疗健康", "AI 医生助理\n病历摘要 · 影像解读\n药物研发辅助", ACCENT_BLUE),
    ("💰 金融科技", "智能投顾\n风控分析 · 合规审查\n自动化报告生成", ACCENT_GREEN),
    ("🏭 智能制造", "生产调度Agent\n预测性维护 · 质量检测\n供应链优化", ORANGE),
    ("🛒 电商零售", "AI 导购助手\n个性化推荐 · 客服\n库存与定价优化", PURPLE),
    ("📚 教育培训", "AI 导师\n自适应学习路径\n作业批改与反馈", RGBColor(0xF7, 0x78, 0x83)),
    ("💻 软件开发", "AI 编程助手\nCode Review · 测试生成\nDevOps 自动化", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(scenes):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.6 + row * 2.8)

    card = add_rect(slide, x, y, Inches(3.7), Inches(2.5), CARD_BG, radius=True)
    # 左侧色条
    add_rect(slide, x, Inches(y.inches + 0.3), Pt(4), Inches(1.9), color)
    add_text_box(slide, Inches(x.inches + 0.3), Inches(y.inches + 0.3),
                 Inches(3.1), Inches(0.5),
                 title, font_size=16, bold=True, color=color)
    add_multiline(slide, Inches(x.inches + 0.3), Inches(y.inches + 0.9),
                  Inches(3.1), Inches(1.5),
                  desc.split("\n"), font_size=12, color=GRAY, line_spacing=1.5)

# ═══════════════════════════════════════════════════
#  第7页: 主流产品对比
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, 6, "主流产品与框架对比", "PRODUCT COMPARISON")

# 表格头部
rows_data = [
    ["产品/框架", "公司", "核心能力", "工具调用", "多Agent", "开源"],
    ["GPTs / Assistant", "OpenAI", "多模态 · Code Interpreter", "✅", "❌", "❌"],
    ["Claude + MCP", "Anthropic", "长上下文 · Computer Use", "✅ MCP", "❌", "❌"],
    ["Gemini", "Google", "多模态 · 长上下文1M", "✅", "❌", "❌"],
    ["Copilot Studio", "Microsoft", "低代码 · M365 集成", "✅", "✅", "❌"],
    ["CrewAI", "开源社区", "角色分工 · 任务编排", "✅", "✅", "✅"],
    ["AutoGen", "Microsoft", "对话式多Agent · 人机协同", "✅", "✅", "✅"],
    ["LangGraph", "LangChain", "有状态图 · 循环推理", "✅", "✅", "✅"],
]

# 简易表格
table_top = Inches(1.6)
col_widths = [Inches(2.3), Inches(1.6), Inches(3.0), Inches(1.5), Inches(1.3), Inches(1.0)]
row_height = Inches(0.58)

for r, row_data in enumerate(rows_data):
    y = table_top + r * row_height
    x_acc = Inches(0.8)
    for c, cell_text in enumerate(row_data):
        w = col_widths[c]
        is_header = (r == 0)
        bg = ACCENT_BLUE if is_header else (CARD_BG if r % 2 == 0 else RGBColor(0x1C, 0x22, 0x2A))
        txt_color = WHITE if is_header else (ACCENT_GREEN if "✅" in cell_text and cell_text != "✅" else
                                              (GRAY if "❌" in cell_text else WHITE))
        rect = add_rect(slide, x_acc, y, w, row_height, bg)
        add_text_box(slide, Inches(x_acc.inches + 0.15), Inches(y.inches + 0.12),
                     Inches(w.inches - 0.3), Inches(0.35),
                     cell_text, font_size=12 if not is_header else 13,
                     bold=is_header, color=txt_color,
                     alignment=PP_ALIGN.CENTER)
        x_acc += w

# ═══════════════════════════════════════════════════
#  第8页: 市场趋势与数据
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, 7, "市场趋势与数据", "MARKET TRENDS & DATA")

# 四个数据卡片
stats = [
    ("$50B+", "2025年 AI Agent\n全球市场规模", ACCENT_BLUE),
    ("47.5%", "CAGR 复合年\n增长率 (2024-2030)", ACCENT_GREEN),
    ("85%", "企业将在2026年\n部署Agent应用", ORANGE),
    ("3x", "Agent化改造后\nROI 提升倍数", PURPLE),
]

for i, (num, desc, color) in enumerate(stats):
    x = Inches(0.8 + i * 3.15)
    card = add_rect(slide, x, Inches(1.6), Inches(2.85), Inches(2.2),
                    CARD_BG, radius=True)
    add_text_box(slide, Inches(x.inches + 0.3), Inches(1.85), Inches(2.3), Inches(0.8),
                 num, font_size=36, bold=True, color=color,
                 alignment=PP_ALIGN.CENTER)
    add_multiline(slide, Inches(x.inches + 0.3), Inches(2.65), Inches(2.3), Inches(0.9),
                  desc.split("\n"), font_size=12, color=GRAY, line_spacing=1.3)

# 趋势
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
             "▎五大趋势", font_size=18, bold=True, color=ACCENT_BLUE)

trends = [
    "1. Agent 从「Copilot」走向「Autopilot」—— 自主决策能力持续增强",
    "2. 多Agent协作成为主流架构 —— 专业化分工、群体智能涌现",
    "3. Agent 与企业系统深度融合 —— ERP/CRM/MES 全面Agent化",
    "4. 小模型 + Agent 崛起 —— 端侧部署、低延迟、隐私合规",
    "5. 监管框架加速成型 —— EU AI Act、中国生成式AI管理办法",
]
add_multiline(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(2.5),
              trends, font_size=14, color=WHITE, line_spacing=1.6)

# ═══════════════════════════════════════════════════
#  第9页: 挑战与伦理
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, 8, "挑战与伦理", "CHALLENGES & ETHICS")

challenges = [
    ("🛡️ 安全风险", "提示注入攻击\n工具滥用与越权\n数据泄露风险", ACCENT_BLUE),
    ("🎭 幻觉问题", "事实性错误生成\n推理偏差放大\n误导性输出", RGBColor(0xF7, 0x78, 0x83)),
    ("⚖️ 责任归属", "Agent决策失误谁担责？\n法律主体地位不明\n监管空白地带", ORANGE),
    ("🌍 社会影响", "就业结构冲击\n算法偏见与公平\n数字鸿沟加剧", PURPLE),
]

for i, (title, desc, color) in enumerate(challenges):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.1)
    y = Inches(1.6 + row * 2.7)

    card = add_rect(slide, x, y, Inches(5.6), Inches(2.4), CARD_BG, radius=True)
    add_rect(slide, x, Inches(y.inches + 0.25), Pt(4), Inches(1.9), color)
    add_text_box(slide, Inches(x.inches + 0.35), Inches(y.inches + 0.25),
                 Inches(5.0), Inches(0.5),
                 title, font_size=18, bold=True, color=color)
    add_multiline(slide, Inches(x.inches + 0.35), Inches(y.inches + 0.85),
                  Inches(5.0), Inches(1.4),
                  desc.split("\n"), font_size=13, color=GRAY, line_spacing=1.6)

# ═══════════════════════════════════════════════════
#  第10页: 封底
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_circle(slide, Inches(10.5), Inches(0.5), Inches(2.5),
           RGBColor(0x1A, 0x3A, 0x5C))
add_circle(slide, Inches(-0.8), Inches(5.5), Inches(1.8),
           RGBColor(0x0D, 0x2B, 0x1A))

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.0),
             "Thank You", font_size=56, bold=True, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.6),
             "AI Agent 正在重新定义人机协作的边界",
             font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.2), Inches(2.3), Pt(2), ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
             "智能化 · 自主化 · 普惠化",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.6), Inches(10), Inches(0.4),
             "© 2025  ·  技术前沿专题  ·  AI Agent 深度介绍",
             font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ── 保存 ──────────────────────────────────────────
output_path = os.path.expanduser("~/Desktop/AI_Agent_深度介绍.pptx")
prs.save(output_path)
print(f"✅ PPT 已生成: {output_path}")
print(f"   共 {len(prs.slides)} 页 | 16:9 宽屏 | 深色科技风")
